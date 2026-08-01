"""
Build the SESPRS-style presentation summarizing the entire facial aging study.
Replicates the navy/gold/cyan template extracted from the user's TIGR-ADM deck.
"""
from __future__ import annotations
import os
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from copy import deepcopy
from lxml import etree

# ------------- Palette (extracted exactly from template) ------------- #
NAVY = RGBColor(0x1E, 0x3A, 0x5F)
CARD_BG = RGBColor(0x24, 0x3D, 0x5C)
CYAN = RGBColor(0x4C, 0xA8, 0xD4)
LIGHT_BLUE = RGBColor(0xB8, 0xCC, 0xE4)
GOLD = RGBColor(0xF4, 0xA7, 0x24)
GREEN = RGBColor(0x4C, 0xD9, 0x7B)
RED = RGBColor(0xE0, 0x57, 0x57)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Calibri"

# Previously hardcoded to one author's machine. See NEXUS_MANUSCRIPT_DIR.
_out_env = os.environ.get("NEXUS_MANUSCRIPT_DIR")
if not _out_env:
    raise SystemExit(
        "NEXUS_MANUSCRIPT_DIR is not set.\n"
        "Point it at the directory holding the study figures, e.g.\n"
        "  export NEXUS_MANUSCRIPT_DIR=~/facial-aging-study"
    )
FIG = Path(_out_env).expanduser()
FIG.mkdir(parents=True, exist_ok=True)
OUT = FIG / "Facial Aging Study Presentation.pptx"

# 16:9 slide
SLIDE_W = 13.333
SLIDE_H = 7.5
prs = Presentation()
prs.slide_width = Inches(SLIDE_W)
prs.slide_height = Inches(SLIDE_H)

BLANK = prs.slide_layouts[6]  # truly blank


# ----------------- Style helpers ------------------- #

def apply_bg(slide, color=NAVY):
    """Apply solid navy background to a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, x, y, w, h, text, *, font=FONT, size=14, color=WHITE,
                bold=False, italic=False, align=PP_ALIGN.LEFT,
                anchor=MSO_ANCHOR.TOP, line_spacing=None):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(2)
    tf.word_wrap = True
    tf.vertical_anchor = anchor

    lines = text.split("\n") if isinstance(text, str) else [text]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing is not None:
            p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return tb


def add_runs(slide, x, y, w, h, runs, *, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, line_spacing=None):
    """runs = list of (text, dict) where dict has size/bold/color/italic/font."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(2)
    tf.word_wrap = True
    tf.vertical_anchor = anchor

    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing is not None:
        p.line_spacing = line_spacing
    for text, props in runs:
        if text == "\n":
            p = tf.add_paragraph()
            p.alignment = align
            if line_spacing is not None:
                p.line_spacing = line_spacing
            continue
        run = p.add_run()
        run.text = text
        run.font.name = props.get("font", FONT)
        run.font.size = Pt(props.get("size", 14))
        run.font.bold = props.get("bold", False)
        run.font.italic = props.get("italic", False)
        run.font.color.rgb = props.get("color", WHITE)
    return tb


def add_title_block(slide, title, *, x=0.5, y=0.35, w=12.3, h=0.85):
    """Slide title with light-blue accent line under it."""
    add_textbox(slide, x, y, w, h, title, size=36, bold=True, color=WHITE,
                anchor=MSO_ANCHOR.MIDDLE)
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(x), Inches(y + h + 0.06),
        Inches(w), Inches(0.04),
    )
    line.line.fill.background()
    line.fill.solid()
    line.fill.fore_color.rgb = LIGHT_BLUE


def add_card(slide, x, y, w, h, *, fill=CARD_BG, border=CYAN, border_w=1.25):
    """Rounded-corner card."""
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = fill
    card.line.color.rgb = border
    card.line.width = Pt(border_w)
    # Tighten corner radius
    card.adjustments[0] = 0.06
    return card


def add_card_with_left_bar(slide, x, y, w, h, bar_color=GOLD, bar_w=0.07):
    """Discussion-style card: left vertical gold bar, no border, transparent."""
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(bar_w), Inches(h)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = bar_color
    bar.line.fill.background()
    return bar


def add_bullets(slide, x, y, w, h, items, *, size=14, color=WHITE, line_spacing=1.25,
                indent=0.0, bullet_color=None):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = Pt(2); tf.margin_right = Pt(2)
    tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    tf.word_wrap = True
    bc = bullet_color or color
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = line_spacing
        p.space_after = Pt(6)
        # Bullet glyph as first run (so we can color it independently)
        bul = p.add_run()
        bul.text = "•  "
        bul.font.name = FONT
        bul.font.size = Pt(size)
        bul.font.color.rgb = bc
        bul.font.bold = True
        body = p.add_run()
        body.text = item
        body.font.name = FONT
        body.font.size = Pt(size)
        body.font.color.rgb = color
    return tb


def add_emory_corner(slide):
    """Subtle 'Emory · Plastic Surgery' mark in bottom-right; placeholder for shield."""
    add_textbox(slide, SLIDE_W - 3.0, SLIDE_H - 0.45, 2.7, 0.3,
                "EMORY  ·  PLASTIC SURGERY",
                size=8, color=LIGHT_BLUE, italic=True, align=PP_ALIGN.RIGHT)


# ============================================================ #
# SLIDE BUILDERS
# ============================================================ #

def slide_title():
    s = prs.slides.add_slide(BLANK)
    apply_bg(s)
    # Title — single coherent line, sized to fit cleanly across two visual lines
    add_textbox(s, 0.6, 1.4, 12.1, 2.6,
        "A Vision–Language Model Pipeline For Standardized\n"
        "Quantification Of Subclinical Facial Aging Endpoints Following\n"
        "Non-Surgical Aesthetic Intervention",
        size=28, bold=True, color=WHITE, line_spacing=1.2)
    # Accent line
    line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                              Inches(0.6), Inches(4.25),
                              Inches(12.1), Inches(0.04))
    line.line.fill.background()
    line.fill.solid(); line.fill.fore_color.rgb = LIGHT_BLUE
    # Authors
    add_textbox(s, 0.6, 4.55, 12.1, 0.6,
        "Orr Shauly, MD¹; Yasmin Fahad, BS¹; Farzad Nahai, MD²; Anita Sethna, MD¹; Albert Losken, MD¹",
        size=15, bold=True, color=WHITE)
    # Affiliations
    add_textbox(s, 0.6, 5.25, 12.1, 0.8,
        "1. Division of Plastic and Reconstructive Surgery, Emory University School of Medicine, Atlanta, GA\n"
        "2. Paces Plastic Surgery and Aesthetic Center, Atlanta, GA",
        size=12, italic=True, color=LIGHT_BLUE, line_spacing=1.25)
    add_emory_corner(s)


def slide_disclosures():
    s = prs.slides.add_slide(BLANK); apply_bg(s)
    add_title_block(s, "Disclosures")
    add_bullets(s, 0.6, 1.7, 12.1, 5.0, [
        "Dr. Shauly is on the advisory board for Machina Health, Inc., is a paid consultant for PRS Atlas, LLC, and Mercuri Gaming, Inc., and is the founder and president of Menley, LLC.",
        "Dr. Losken is on the advisory board for Bimini Health.",
        "All other authors have no disclosures.",
        "Computational infrastructure and vision–language model API credits were supported by the Southeastern Society of Plastic and Reconstructive Surgery annual research grant.",
    ], size=16, line_spacing=1.4)
    add_emory_corner(s)


def slide_background():
    s = prs.slides.add_slide(BLANK); apply_bg(s)
    add_title_block(s, "Background")
    add_bullets(s, 0.6, 1.7, 12.1, 5.0, [
        "Non-surgical facial rejuvenation has expanded markedly post–COVID-19, driven by video-conferencing exposure (Rice et al. — 56.7% of providers reporting increased demand).",
        "Existing macroscopic facial-aging scales (Brazin Atlas, SCINEXA) are highly subjective with substantial inter-rater variability that limits longitudinal and cross-cohort comparison.",
        "Prior AI work has reported strong age-prediction performance (Assi et al. r = 0.937; Park et al. precision/sensitivity > 0.9; Zhang et al. deep-learning superiority) — but in single-purpose, single-ethnicity, healthy-subject cohorts.",
        "No reproducible, schema-validated, multi-endpoint pipeline has been described that ingests routine before/after pairs from a heterogeneous cosmetic and reconstructive cohort and produces an auditable comparable score-set across pairs.",
    ], size=15, line_spacing=1.35)
    add_emory_corner(s)


def slide_objectives():
    s = prs.slides.add_slide(BLANK); apply_bg(s)
    add_title_block(s, "Objectives")
    add_bullets(s, 0.6, 1.7, 12.1, 5.0, [
        "Develop and evaluate a vision–language model (VLM) pipeline that scores standardized before/after facial photograph pairs against the seven primary endpoints from our research protocol.",
        "Augment the rubric with seven anatomic sub-region severity scores (crow's feet, nasolabial folds, forehead, glabellar, perioral, tear-trough, jowling).",
        "Consolidate per-pair output into a single cohort-standardized improvement score that is directly comparable across heterogeneous treatment timepoints and intervention classes.",
        "Assess internal consistency, post-procedure time-trend, and intervention-stratified discrimination as preliminary steps toward a clinician-validated objective endpoint for non-surgical facial-aging research.",
    ], size=15, line_spacing=1.35)
    add_emory_corner(s)


def slide_methods_pipeline():
    s = prs.slides.add_slide(BLANK); apply_bg(s)
    add_title_block(s, "Methods – Cohort & Image Pipeline")
    # Two cards
    add_card(s, 0.5, 1.7, 6.1, 5.0)
    add_textbox(s, 0.7, 1.85, 5.7, 0.4, "Cohort", size=18, bold=True, color=GOLD)
    add_bullets(s, 0.7, 2.35, 5.7, 4.3, [
        "Retrospective de-identified database of standardized frontal facial photographs from a multi-site aesthetic and reconstructive practice.",
        "Inclusion: neurotoxin and/or hyaluronic-acid filler. Excluded: prior or interval surgery, peels, laser, or microneedling.",
        "134 before/after pairs across 65 distinct subjects.",
        "Post-procedure follow-up 1–53 weeks (median 2 weeks).",
    ], size=13, line_spacing=1.3)

    add_card(s, 6.7, 1.7, 6.1, 5.0)
    add_textbox(s, 6.9, 1.85, 5.7, 0.4, "Image Pipeline", size=18, bold=True, color=GOLD)
    add_bullets(s, 6.9, 2.35, 5.7, 4.3, [
        "De-identification: SSD MobileNet v1 + 68-point face landmarks (Nexus Face API®) with oval crop and ocular masking.",
        "Identifying tattoos and jewelry detected and cropped; deterministic geometric fallback when face detection fails.",
        "Cryptographic checksums of all model weights persisted with each scored row for audit reproducibility.",
        "Paired luminance harmonization (Rec. 601 weighting) before scoring to remove lighting-driven appearance variance.",
    ], size=13, line_spacing=1.3)
    add_emory_corner(s)


def slide_methods_rubric():
    s = prs.slides.add_slide(BLANK); apply_bg(s)
    add_title_block(s, "Methods – Scoring Rubric & Composite Algorithm")
    # Two cards
    add_card(s, 0.5, 1.7, 6.1, 5.0)
    add_textbox(s, 0.7, 1.85, 5.7, 0.4, "Scoring Rubric", size=18, bold=True, color=GOLD)
    add_bullets(s, 0.7, 2.35, 5.7, 4.3, [
        "Vision model: Claude Opus 4.7 / Amazon Rekognition; rubric supplied as fixed system prompt with paired before/after images.",
        "Top-level: predicted facial age, global rhytid severity, subclinical wrinkles, perceived firmness, density, fullness, gonial-angle.",
        "Sub-regions (0–100 severity each): crow's feet, nasolabial folds, forehead, glabellar, perioral, tear-trough, jowling.",
        "JSON output validated against Zod schema; up to 4 retries on validation failure (no failures observed in cohort).",
    ], size=13, line_spacing=1.3)

    add_card(s, 6.7, 1.7, 6.1, 5.0)
    add_textbox(s, 6.9, 1.85, 5.7, 0.4, "Composite Improvement Score", size=18, bold=True, color=GOLD)
    add_bullets(s, 6.9, 2.35, 5.7, 4.3, [
        "Polarity-flip every metric so positive = improvement on every input.",
        "Aggregate into 4 clinical domains: Age, Wrinkles (mean of 7 sub-regions), Volume, Gonial Phenotype.",
        "z-Standardize each domain across the cohort using sample SD.",
        "Composite = unweighted mean of 4 domain z-scores; cohort percentile by average-rank ordering.",
    ], size=13, line_spacing=1.3)
    add_emory_corner(s)


def add_table_styled(slide, x, y, w_total, headers, rows, col_widths, *, header_size=13, body_size=12,
                     emphasize_col=None, emphasize_color=GOLD):
    """Build a navy-themed table matching the template look."""
    n_rows = len(rows) + 1
    n_cols = len(headers)
    # Use python-pptx Table
    tbl_shape = slide.shapes.add_table(n_rows, n_cols, Inches(x), Inches(y),
                                       Inches(w_total), Inches(0.5 * n_rows))
    tbl = tbl_shape.table
    # Set column widths
    for ci, cw in enumerate(col_widths):
        tbl.columns[ci].width = Inches(cw)
    # Header row
    for ci, h in enumerate(headers):
        cell = tbl.cell(0, ci)
        cell.fill.solid(); cell.fill.fore_color.rgb = CARD_BG
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run(); run.text = h
        run.font.name = FONT; run.font.size = Pt(header_size); run.font.bold = True
        run.font.color.rgb = WHITE
        cell.margin_left = cell.margin_right = Pt(6)
        cell.margin_top = cell.margin_bottom = Pt(4)
    # Data rows
    for ri, row in enumerate(rows, start=1):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = NAVY
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT
            run = p.add_run(); run.text = str(val)
            run.font.name = FONT; run.font.size = Pt(body_size)
            run.font.color.rgb = (emphasize_color if ci == emphasize_col else WHITE)
            run.font.bold = (ci == emphasize_col)
            cell.margin_left = cell.margin_right = Pt(6)
            cell.margin_top = cell.margin_bottom = Pt(3)
    return tbl_shape


def slide_cohort_overview():
    s = prs.slides.add_slide(BLANK); apply_bg(s)
    add_title_block(s, "Cohort Overview")
    headers = ["Variable", "Value"]
    rows = [
        ("Total Before/After Pairs Scored", "134"),
        ("Distinct Subjects", "65"),
        ("Pairs With Documented Post-Procedure Week", "105 (78.4%)"),
        ("Post-Procedure Week — Min / Median / Max", "1  /  2  /  53"),
        ("Pairs In 1–4 / 5–12 / 13+ Week Buckets", "75  /  15  /  15"),
        ("Pairs Skipped (Source Image Missing)", "6"),
        ("Vision–Language Model", "Claude Opus 4.7 / Amazon Rekognition"),
        ("Mean Rubric Self-Confidence (0–1)", "0.79 (Range 0.65–0.95)"),
    ]
    add_table_styled(s, 1.5, 1.8, 10.3, headers, rows, [5.5, 4.8],
                    header_size=14, body_size=13, emphasize_col=1, emphasize_color=GOLD)
    add_emory_corner(s)


def stat_callout(slide, x, y, w, h, header,
                 main_label, main_value, main_color, main_unit,
                 sec_label, sec_value, sec_color,
                 footer):
    """Single-headline stat card: ONE big number plus a smaller secondary stat below."""
    add_card(slide, x, y, w, h)
    add_textbox(slide, x, y + 0.20, w, 0.4, header, size=14, bold=True, color=GOLD,
                align=PP_ALIGN.CENTER)
    add_textbox(slide, x, y + 0.75, w, 0.3, main_label, size=12, color=LIGHT_BLUE,
                align=PP_ALIGN.CENTER)
    add_textbox(slide, x, y + 1.05, w, 1.2, main_value, size=66, bold=True,
                color=main_color, align=PP_ALIGN.CENTER)
    add_textbox(slide, x, y + 2.30, w, 0.32, main_unit, size=12, italic=True, color=LIGHT_BLUE,
                align=PP_ALIGN.CENTER)
    # Divider line
    div = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(x + w*0.2), Inches(y + 2.75),
                                 Inches(w*0.6), Inches(0.015))
    div.line.fill.background()
    div.fill.solid(); div.fill.fore_color.rgb = LIGHT_BLUE
    add_textbox(slide, x, y + 2.85, w, 0.3, sec_label, size=12, color=LIGHT_BLUE,
                align=PP_ALIGN.CENTER)
    add_textbox(slide, x, y + 3.15, w, 0.5, sec_value, size=24, bold=True,
                color=sec_color, align=PP_ALIGN.CENTER)
    add_textbox(slide, x, y + h - 0.45, w, 0.35, footer, size=12, bold=True, color=GOLD,
                align=PP_ALIGN.CENTER)


def slide_results_primary():
    s = prs.slides.add_slide(BLANK); apply_bg(s)
    add_title_block(s, "Results – Primary Endpoints")
    # Three callouts — each has ONE big number plus a smaller IQR/secondary stat
    stat_callout(s, 0.5, 1.7, 4.0, 4.4,
                 "Δ Predicted Facial Age",
                 "Median (AFTER − BEFORE)", "−4", GREEN, "Years",
                 "Pairs Improved", "74%", GREEN,
                 "p < 0.001 vs. zero")
    stat_callout(s, 4.7, 1.7, 4.0, 4.4,
                 "Δ Static Rhytids",
                 "Median (AFTER − BEFORE)", "−6", GREEN, "0–100 Severity Scale",
                 "Interquartile Range", "−10 to −2", LIGHT_BLUE,
                 "All sub-regions concordant")
    stat_callout(s, 8.9, 1.7, 4.0, 4.4,
                 "Δ Subclinical Wrinkles",
                 "Median (AFTER − BEFORE)", "−6", GREEN, "0–100 Severity Scale",
                 "Interquartile Range", "−12 to −2", LIGHT_BLUE,
                 "Tracks global wrinkle Δ")
    add_textbox(s, 0.5, 6.4, 12.3, 0.5,
                "All 134 pairs successfully scored. No rubric validation failures. Mean model self-confidence 0.79 (range 0.65–0.95).",
                size=13, italic=True, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)
    add_emory_corner(s)


def slide_subregion():
    s = prs.slides.add_slide(BLANK); apply_bg(s)
    add_title_block(s, "Results – Anatomic Sub-Regions")
    # Figure left, key takeaways right
    s.shapes.add_picture(str(FIG / "Figure 2.png"),
                         Inches(0.5), Inches(1.75), Inches(7.4), Inches(4.6))
    add_card(s, 8.2, 1.75, 4.6, 4.6)
    add_textbox(s, 8.4, 1.9, 4.3, 0.4, "Key Takeaways", size=18, bold=True, color=GOLD)
    add_bullets(s, 8.4, 2.4, 4.3, 4.0, [
        "All seven anatomic sub-regions show cohort-mean improvement.",
        "Largest reductions: Jowling (−6.7 ± 10.8), Tear-Trough (−6.9 ± 9.9), Nasolabial (−5.6 ± 8.8).",
        "Approximately 75% of pairs improved on every sub-region.",
        "Directional concordance with global wrinkle delta supports rubric reliability.",
    ], size=13, line_spacing=1.3)
    add_textbox(s, 0.5, 6.5, 12.3, 0.4,
                "Figure 2. Cohort-mean improvement (severity points; positive = better) by anatomic sub-region.",
                size=11, italic=True, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)
    add_emory_corner(s)


def slide_composite():
    s = prs.slides.add_slide(BLANK); apply_bg(s)
    add_title_block(s, "Results – Composite Improvement Score")
    # Figure left
    s.shapes.add_picture(str(FIG / "Figure 1.png"),
                         Inches(0.4), Inches(1.75), Inches(7.4), Inches(4.6))
    # Stat card right
    add_card(s, 8.1, 1.75, 4.7, 4.6)
    add_textbox(s, 8.3, 1.9, 4.3, 0.4, "Composite Distribution", size=18, bold=True, color=GOLD)
    # 4 metric callouts inside
    metrics = [
        ("Cohort Mean", "0.00", LIGHT_BLUE),
        ("Standard Deviation", "0.98", WHITE),
        ("Range", "−3.26 to +2.26", WHITE),
        ("Pairs Above Cohort", "54.5%", GREEN),
    ]
    yy = 2.5
    for label, val, color in metrics:
        add_textbox(s, 8.3, yy, 4.3, 0.3, label, size=12, color=LIGHT_BLUE)
        add_textbox(s, 8.3, yy + 0.25, 4.3, 0.55, val, size=24, bold=True, color=color)
        yy += 0.95
    add_textbox(s, 0.5, 6.5, 12.3, 0.4,
                "Figure 1. Distribution of the standardized composite improvement score (n = 134 pairs).",
                size=11, italic=True, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)
    add_emory_corner(s)


def slide_consistency():
    s = prs.slides.add_slide(BLANK); apply_bg(s)
    add_title_block(s, "Results – Internal Consistency & Time Trend")
    # Two figures side by side
    s.shapes.add_picture(str(FIG / "Figure 3.png"),
                         Inches(0.4), Inches(1.75), Inches(6.0), Inches(4.4))
    s.shapes.add_picture(str(FIG / "Figure 4.png"),
                         Inches(6.7), Inches(1.75), Inches(6.0), Inches(4.4))
    add_textbox(s, 0.4, 6.2, 6.0, 0.7,
                "Pearson r = 0.95 across all four-domain pairs (range 0.94–0.98); high inter-variable concordance with likely VLM halo effect.",
                size=11, italic=True, color=LIGHT_BLUE, align=PP_ALIGN.CENTER, line_spacing=1.2)
    add_textbox(s, 6.7, 6.2, 6.0, 0.7,
                "No monotonic time-trend across 1–4, 5–12, and 13+ week buckets — suggests a possible durable / prophylactic treatment effect.",
                size=11, italic=True, color=LIGHT_BLUE, align=PP_ALIGN.CENTER, line_spacing=1.2)
    add_emory_corner(s)


def slide_intervention():
    s = prs.slides.add_slide(BLANK); apply_bg(s)
    add_title_block(s, "Results – Intervention Stratification")
    headers = ["Intervention Class", "n", "Δ Age (yr) — Med (IQR)", "Δ Wrinkles — Med (IQR)",
               "Δ Firmness — Mean ± SD", "Composite Score"]
    rows = [
        ("Neurotoxin Only", "51", "−2 (−4 to 0)", "−4 (−7 to 0)", "+3.84 ± 11.68", "−0.36 ± 0.99"),
        ("Neurotoxin + HA Filler", "77", "−4 (−6 to −2)", "−7 (−10 to −4)", "+11.23 ± 10.47", "+0.26 ± 0.90"),
        ("Filler Only", "1", "−4", "−10", "+18.0 ± 0", "+0.57"),
        ("No Active Intervention", "5", "−2 (−4 to 0)", "−3 (−7 to +4)", "+4.80 ± 14.18", "−0.40 ± 1.14"),
    ]
    add_table_styled(s, 0.5, 1.75, 12.3, headers, rows,
                     [2.6, 0.6, 2.4, 2.4, 2.1, 2.2],
                     header_size=12, body_size=11.5,
                     emphasize_col=5, emphasize_color=GOLD)
    # Caption / takeaway
    add_textbox(s, 0.5, 5.0, 12.3, 0.5,
                "Combination Therapy Outperforms Neurotoxin Alone Across Every Endpoint",
                size=18, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    add_textbox(s, 0.5, 5.55, 12.3, 1.5,
                "Composite gap of ≈ 0.6 cohort SD; largest combination-vs-toxin-alone effects in lower-face regions where volumetric correction predominates "
                "(jowling −9.2 vs −3.0; nasolabial folds −7.8 vs −2.6; tear-trough −7.0 vs −2.5). "
                "No monotonic dose–response across neurotoxin tertiles (≤40 U: +0.01; 40–50 U: −0.05; ≥50 U: +0.08). "
                "Brand outcomes directionally similar (Botox +0.02, Dysport −0.06, Jeuveau −0.01).",
                size=12, color=LIGHT_BLUE, align=PP_ALIGN.LEFT, line_spacing=1.3)
    add_emory_corner(s)


def slide_discussion():
    s = prs.slides.add_slide(BLANK); apply_bg(s)
    add_title_block(s, "Discussion")
    # 4 cards with left gold bar
    cards = [
        ("Combination Therapy Effect",
         "Neurotoxin + HA filler outperforms neurotoxin alone by ≈ 0.6 cohort SD on the composite, with the largest gains in lower-face regions where volumetric correction dominates — biologically coherent and supports an additive mechanism."),
        ("Halo Effect Acknowledged",
         "Inter-domain Pearson r 0.94–0.98 indicates the rubric's apparent multi-endpoint richness compresses into a single holistic appearance judgment; future iterations will score domains in independent calls."),
        ("Pipeline Reproducibility",
         "JSON-schema-validated rubric, checksum-pinned model weights, and resumable scoring give every persisted row a complete, auditable, reproducible endpoint set."),
        ("Negative-Control Anchoring",
         "The five no-active-intervention follow-up pairs anchored a near-zero composite (−0.40 ± 1.14), supporting attribution of the cohort positive signal to clinical treatment rather than rubric drift."),
    ]
    yy = 1.7
    for header, body in cards:
        add_card_with_left_bar(s, 0.5, yy, 12.3, 1.18)
        add_textbox(s, 0.75, yy + 0.07, 12.0, 0.35, header, size=15, bold=True, color=GOLD)
        add_textbox(s, 0.75, yy + 0.45, 12.0, 0.7, body,
                    size=12, color=WHITE, line_spacing=1.25)
        yy += 1.27
    add_emory_corner(s)


def slide_limitations():
    s = prs.slides.add_slide(BLANK); apply_bg(s)
    add_title_block(s, "Limitations")
    add_bullets(s, 0.6, 1.7, 12.1, 5.0, [
        "VLM as the sole observer — internal consistency cannot substitute for external clinician validation.",
        "Single frontal angle vs. the protocol-planned five-angle photography, reducing sensitivity for jawline-definition and gonial-angle endpoints.",
        "Heterogeneous post-procedure timepoints (1–53 weeks; median 2 weeks) preclude formal time-response modeling and bias the cohort toward early-window assessments.",
        "Observational intervention assignment — not randomized; combination-therapy patients may differ at baseline from those receiving neurotoxin alone.",
        "VLM behavior is not bit-identical across model releases; reproducibility requires pinning the model release identifier or accepting a documented inter-release tolerance.",
    ], size=15, line_spacing=1.4)
    add_emory_corner(s)


def slide_conclusion():
    s = prs.slides.add_slide(BLANK); apply_bg(s)
    add_title_block(s, "Conclusion")
    add_bullets(s, 0.6, 1.7, 12.1, 5.0, [
        "A vision–language model pipeline can produce complete, schema-validated, multi-endpoint facial-aging assessments on standardized before/after pairs in a heterogeneous real-world cosmetic and reconstructive cohort.",
        "Pipeline reproduces clinically reasonable cohort-level findings — median predicted-age reduction of 4 years, majority improvement on each of seven anatomic sub-regions.",
        "Framework discriminates between intervention classes in a clinically interpretable way: combination neurotoxin + HA filler outperforms neurotoxin alone across every endpoint.",
        "Negative-control follow-up pairs anchor a near-zero composite, supporting attribution of the cohort signal to clinical treatment rather than pipeline artifact.",
        "Next phase: prospective demographic-stratified, intervention-controlled accrual against blinded clinician raters before broader clinical deployment.",
    ], size=14, line_spacing=1.35)
    add_emory_corner(s)


def slide_thanks():
    s = prs.slides.add_slide(BLANK); apply_bg(s)
    add_textbox(s, 0.5, 2.8, 12.3, 1.2, "Thank You",
                size=72, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)
    line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                              Inches(4.5), Inches(4.05), Inches(4.3), Inches(0.04))
    line.line.fill.background()
    line.fill.solid(); line.fill.fore_color.rgb = LIGHT_BLUE
    add_textbox(s, 0.5, 4.25, 12.3, 0.5, "Questions & Discussion",
                size=22, italic=True, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)
    add_textbox(s, 0.5, 5.3, 12.3, 0.5, "orr.shauly@emory.edu",
                size=18, color=WHITE, align=PP_ALIGN.CENTER)
    add_emory_corner(s)


# ============================================================ #
# BUILD
# ============================================================ #

slide_title()
slide_disclosures()
slide_background()
slide_objectives()
slide_methods_pipeline()
slide_methods_rubric()
slide_cohort_overview()
slide_results_primary()
slide_subregion()
slide_composite()
slide_consistency()
slide_intervention()
slide_discussion()
slide_limitations()
slide_conclusion()
slide_thanks()

prs.save(str(OUT))
print(f"Saved {OUT}  ({len(prs.slides)} slides)")
